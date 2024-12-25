import json
from pptx import Presentation
import re

def replace_text(text) -> str:
    text = text.replace('（', '(');
    text = text.replace('）', ')');
    text = text.replace('：', ':');
    text = text.replace('，', ', ');
    text = text.replace('；', '; ');

    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r' +', ' ', text)

    return text;


# 提取PPT中的文本信息
def extract_text_from_ppt(pptx_file):
    prs = Presentation(pptx_file)
    slides_text = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = ""
        for shape in slide.shapes:
            # 如果形状包含文本框
            if hasattr(shape, "text"):
                slide_text += shape.text.strip() + "\n"
            
        # 如果该幻灯片有文本，则将其添加到列表中
        if slide_text.strip():
            slide_text = replace_text(slide_text);
            slides_text.append({
                "Type": 0,
                "Text": slide_text.strip()
            })

    return slides_text

# 将提取的文本信息保存为JSON格式的txt文件
def save_as_json(slides_text, output_file):
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(slides_text, f, ensure_ascii=False, indent=4)

# 主程序
if __name__ == "__main__":
    pptx_file = "C:\\Users\\Aite\\Desktop\\课件\\第3章  编码和调制.pptx"  # 替换为路径
    output_file = "output.json"       # 输出文件路径

    slides_text = extract_text_from_ppt(pptx_file)
    save_as_json(slides_text, output_file)

    print(f"文本信息已提取并保存到 {output_file}")
